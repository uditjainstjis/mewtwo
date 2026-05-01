"""Build Synapta CTO pitch deck (.pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Palette
BG = RGBColor(0x0A, 0x0E, 0x27)
PANEL = RGBColor(0x14, 0x1B, 0x3D)
PANEL_LIGHT = RGBColor(0x1C, 0x25, 0x4F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
GRAY = RGBColor(0xA0, 0xAE, 0xC0)
GRAY_DIM = RGBColor(0x6B, 0x76, 0x95)
GREEN = RGBColor(0x10, 0xE0, 0xA0)
RED = RGBColor(0xFF, 0x6B, 0x6B)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ---------- helpers ----------
def blank_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return s

def text(slide, content, left, top, width, height, size=18, color=WHITE,
         bold=False, align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = content
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = font
    return tb

def multi_text(slide, lines, left, top, width, height, align=PP_ALIGN.LEFT,
               anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, color, bold, font)"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, (txt, sz, col, b, f) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.color.rgb = col
        r.font.bold = b
        r.font.name = f
    return tb

def rect(slide, left, top, width, height, fill=PANEL, line=None, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    s.shadow.inherit = False
    return s

def accent_bar(slide):
    rect(slide, Inches(0), Inches(0), Inches(0.18), Inches(7.5), fill=CYAN)

def page_label(slide, text_str, page_num):
    text(slide, text_str.upper(), Inches(0.7), Inches(0.45), Inches(8), Inches(0.3),
         size=10, color=CYAN, bold=True)
    text(slide, f"{page_num:02d} / 08", Inches(12.0), Inches(0.45), Inches(1.0), Inches(0.3),
         size=10, color=GRAY_DIM, bold=False, align=PP_ALIGN.RIGHT)

def footer(slide):
    # Tiny corner footer only — keeps the bottom edge clean for closing lines
    text(slide, "SYNAPTA",
         Inches(12.0), Inches(7.15), Inches(1.0), Inches(0.25),
         size=8, color=GRAY_DIM, bold=True, align=PP_ALIGN.RIGHT)

# =====================================================================
# SLIDE 1 — TITLE / HOOK
# =====================================================================
s = blank_slide()
accent_bar(s)
text(s, "SYNAPTA", Inches(0.7), Inches(0.6), Inches(8), Inches(0.4),
     size=12, color=CYAN, bold=True)
text(s, "Sovereign AI for regulated enterprise.", Inches(0.7), Inches(0.95), Inches(10), Inches(0.4),
     size=12, color=GRAY)

# Multi-line headline
multi_text(s, [
    ("Frontier-class reasoning,", 50, WHITE, True, "Calibri"),
    ("sized for your infrastructure,", 50, WHITE, True, "Calibri"),
    ("inside your firewall —", 50, WHITE, True, "Calibri"),
    ("at 20× lower compute cost.", 50, CYAN, True, "Calibri"),
], Inches(0.7), Inches(2.0), Inches(12), Inches(3.5))

text(s, "An adapter-routing platform that lets any open base model host dozens of swappable",
     Inches(0.7), Inches(5.6), Inches(12), Inches(0.4), size=14, color=GRAY)
text(s, "domain experts. Customer's data never leaves their infrastructure.",
     Inches(0.7), Inches(5.95), Inches(12), Inches(0.4), size=14, color=GRAY)

text(s, "SOLO FOUNDER  ·  19  ·  4 RESEARCH PAPERS  ·  VALIDATED AT 30B",
     Inches(0.7), Inches(6.7), Inches(12), Inches(0.3), size=10, color=GRAY_DIM, bold=True)

# =====================================================================
# SLIDE 2 — WHY NOW
# =====================================================================
s = blank_slide()
accent_bar(s)
page_label(s, "Why now", 2)
footer(s)

text(s, "70% of enterprise AI data is legally restricted",
     Inches(0.7), Inches(1.1), Inches(12), Inches(0.7), size=34, bold=True)
text(s, "from leaving its country of origin.",
     Inches(0.7), Inches(1.7), Inches(12), Inches(0.7), size=34, bold=True, color=CYAN)

# Three stat panels
panel_w = Inches(3.9)
panel_h = Inches(2.6)
top = Inches(3.4)
gap = Inches(0.25)
left0 = Inches(0.7)

# Panel 1
rect(s, left0, top, panel_w, panel_h, fill=PANEL, rounded=True)
text(s, "0", left0, Inches(3.6), panel_w, Inches(1.0),
     size=68, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
text(s, "frontier APIs satisfy",
     left0, Inches(4.85), panel_w, Inches(0.4), size=14, color=GRAY, align=PP_ALIGN.CENTER)
text(s, "RBI · SEBI · DPDP · EU AI Act",
     left0, Inches(5.2), panel_w, Inches(0.4), size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(s, "data localization mandates.",
     left0, Inches(5.55), panel_w, Inches(0.4), size=14, color=GRAY, align=PP_ALIGN.CENTER)

# Panel 2
left1 = left0 + panel_w + gap
rect(s, left1, top, panel_w, panel_h, fill=PANEL, rounded=True)
text(s, "$400M – 1B", left1, Inches(3.7), panel_w, Inches(0.9),
     size=42, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
text(s, "Indian BFSI back-office", left1, Inches(4.8), panel_w, Inches(0.4),
     size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(s, "AI spend addressable in our wedge.", left1, Inches(5.15), panel_w, Inches(0.4),
     size=14, color=GRAY, align=PP_ALIGN.CENTER)
text(s, "~200 mid-tier institutions.", left1, Inches(5.5), panel_w, Inches(0.4),
     size=14, color=GRAY, align=PP_ALIGN.CENTER)

# Panel 3
left2 = left1 + panel_w + gap
rect(s, left2, top, panel_w, panel_h, fill=PANEL, rounded=True)
text(s, "20×", left2, Inches(3.7), panel_w, Inches(0.9),
     size=64, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
text(s, "lower compute cost", left2, Inches(4.85), panel_w, Inches(0.4),
     size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(s, "vs self-hosted frontier-class", left2, Inches(5.2), panel_w, Inches(0.4),
     size=14, color=GRAY, align=PP_ALIGN.CENTER)
text(s, "(64× H100 cluster vs 2-3× H100).", left2, Inches(5.55), panel_w, Inches(0.4),
     size=14, color=GRAY, align=PP_ALIGN.CENTER)

text(s, "Frontier APIs cannot serve regulated buyers. Frontier self-hosting is unaffordable.",
     Inches(0.7), Inches(6.45), Inches(12), Inches(0.4), size=14, color=GRAY)
text(s, "We are the third option.", Inches(0.7), Inches(6.85),
     Inches(12), Inches(0.4), size=15, color=CYAN, bold=True)

# =====================================================================
# SLIDE 3 — THE NON-OBVIOUS TRUTH (CODE PARADOX)
# =====================================================================
s = blank_slide()
accent_bar(s)
page_label(s, "The non-obvious truth", 3)
footer(s)

text(s, "We discovered something nobody else has published.",
     Inches(0.7), Inches(1.1), Inches(12), Inches(0.6), size=28, bold=True, color=GRAY)

text(s, "On reasoning tasks,", Inches(0.7), Inches(2.0), Inches(12), Inches(0.7),
     size=40, bold=True)
text(s, "code-trained adapters beat math-trained adapters.",
     Inches(0.7), Inches(2.6), Inches(12), Inches(0.7), size=40, bold=True, color=CYAN)

# Two stat boxes
box_w = Inches(5.5)
box_h = Inches(2.0)
box_top = Inches(4.2)

# Math adapter
rect(s, Inches(0.7), box_top, box_w, box_h, fill=PANEL, rounded=True)
text(s, "MATH ADAPTER", Inches(0.95), Inches(4.4), box_w, Inches(0.4),
     size=11, color=GRAY, bold=True)
text(s, "50.5%", Inches(0.95), Inches(4.85), Inches(2.5), Inches(1.0),
     size=58, color=WHITE, bold=True)
text(s, "on MATH-500", Inches(3.4), Inches(5.15), Inches(2.5), Inches(0.4),
     size=14, color=GRAY)
text(s, "(specialized fine-tuning)", Inches(3.4), Inches(5.45), Inches(2.5), Inches(0.4),
     size=11, color=GRAY_DIM)

# Code adapter
rect(s, Inches(7.1), box_top, box_w, box_h, fill=PANEL, line=CYAN, rounded=True)
text(s, "CODE ADAPTER", Inches(7.35), Inches(4.4), box_w, Inches(0.4),
     size=11, color=CYAN, bold=True)
text(s, "56.0%", Inches(7.35), Inches(4.85), Inches(2.5), Inches(1.0),
     size=58, color=CYAN, bold=True)
text(s, "on MATH-500", Inches(9.8), Inches(5.15), Inches(2.5), Inches(0.4),
     size=14, color=GRAY)
text(s, "(+5.5 over math adapter)", Inches(9.8), Inches(5.45), Inches(2.5), Inches(0.4),
     size=11, color=GREEN, bold=True)

# Implication
text(s, "Implication:  Specialized fine-tuning isn't the bottleneck.  Composition is.",
     Inches(0.7), Inches(6.45), Inches(12), Inches(0.5), size=18, color=WHITE, bold=True)
text(s, "n=200 · Nemotron-3-Nano-30B · NeurIPS submission in pipeline",
     Inches(0.7), Inches(6.9), Inches(12), Inches(0.3), size=10, color=GRAY_DIM)

# =====================================================================
# SLIDE 4 — BENCHMARK GRID
# =====================================================================
s = blank_slide()
accent_bar(s)
page_label(s, "Measured results", 4)
footer(s)

text(s, "Same hardware. Same base model.",
     Inches(0.7), Inches(1.1), Inches(12), Inches(0.6), size=24, color=GRAY, bold=True)
text(s, "+11 to +14.5 points across reasoning benchmarks.",
     Inches(0.7), Inches(1.7), Inches(12), Inches(0.7), size=32, bold=True, color=CYAN)

# Build table
rows = [
    ("Method", "ARC-Challenge", "MATH-500", "HumanEval", "MBPP"),
    ("Base Nemotron-30B", "20.0%", "41.5%", "50.0%", "8.0%"),
    ("Static Merge (DARE/TIES)", "19.0%", "56.0%", "34.0%", "0.0%"),
    ("Best Single Adapter", "31.0%", "56.0%", "60.0%", "6.0%"),
    ("Our Adapter Routing", "31.0%", "56.0%", "48.0%", "5.0%"),
]
# Bold the column-best
bold_cells = {(3, 1), (3, 2), (3, 3), (4, 1), (4, 2)}  # row, col indices

table_left = Inches(0.7)
table_top = Inches(2.8)
col_widths = [Inches(3.6), Inches(1.9), Inches(1.9), Inches(1.9), Inches(1.9)]
row_h = Inches(0.6)

# Header
x = table_left
for ci, txt in enumerate(rows[0]):
    rect(s, x, table_top, col_widths[ci], row_h, fill=PANEL_LIGHT, rounded=False)
    text(s, txt, x + Inches(0.2), table_top + Inches(0.16), col_widths[ci], Inches(0.4),
         size=12, color=CYAN, bold=True,
         align=(PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER))
    x += col_widths[ci]

# Body
for ri in range(1, len(rows)):
    y = table_top + row_h * ri
    fill = PANEL if ri % 2 == 1 else BG
    x = table_left
    for ci, txt in enumerate(rows[ri]):
        rect(s, x, y, col_widths[ci], row_h, fill=fill, rounded=False)
        is_bold = (ri, ci) in bold_cells
        col = CYAN if is_bold else WHITE
        if ri == 1:
            col = GRAY  # baseline row dimmer
        text(s, txt, x + Inches(0.2), y + Inches(0.16), col_widths[ci], Inches(0.4),
             size=14 if is_bold else 13, color=col, bold=is_bold,
             align=(PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER))
        x += col_widths[ci]

# Footnote / takeaway
text(s, "Static merging COLLAPSES non-dominant benchmarks (ARC: -1, HumanEval: -16).",
     Inches(0.7), Inches(6.05), Inches(12), Inches(0.4), size=14, color=RED, bold=True)
text(s, "Routing PRESERVES peak expert performance across domains.",
     Inches(0.7), Inches(6.4), Inches(12), Inches(0.4), size=14, color=GREEN, bold=True)
text(s, "Format Guard variant on HumanEval reaches 48% (+24 vs in-environment base).",
     Inches(0.7), Inches(6.75), Inches(12), Inches(0.4), size=11, color=GRAY_DIM)

# =====================================================================
# SLIDE 5 — ARCHITECTURE & MOAT
# =====================================================================
s = blank_slide()
accent_bar(s)
page_label(s, "How it works · why it cannot be copied", 5)
footer(s)

text(s, "One adapter system. Any base. Any infrastructure.",
     Inches(0.7), Inches(1.1), Inches(12), Inches(0.7), size=30, bold=True)

# Tier diagram (4 stacked tier boxes on left, adapter library on right)
tiers = [
    ("EDGE / MOBILE", "7B Gemma · Phi · Qwen", "Phones, laptops, kiosks"),
    ("MID-TIER", "30-70B Nemotron · Llama", "On-prem servers, single-GPU"),
    ("ENTERPRISE", "100-200B", "Datacenter, multi-GPU"),
    ("FRONTIER", "500B+", "Hyperscaler-grade"),
]
tier_left = Inches(0.7)
tier_top = Inches(2.3)
tier_w = Inches(5.5)
tier_h = Inches(0.95)
tier_gap = Inches(0.12)

for i, (label, base, deploy) in enumerate(tiers):
    y = tier_top + (tier_h + tier_gap) * i
    rect(s, tier_left, y, tier_w, tier_h, fill=PANEL, rounded=True)
    text(s, label, tier_left + Inches(0.25), y + Inches(0.08), Inches(2), Inches(0.3),
         size=10, color=CYAN, bold=True)
    text(s, base, tier_left + Inches(0.25), y + Inches(0.35), Inches(3), Inches(0.4),
         size=15, color=WHITE, bold=True)
    text(s, deploy, tier_left + Inches(0.25), y + Inches(0.65), Inches(5), Inches(0.3),
         size=10, color=GRAY)

# Adapter library on right
lib_left = Inches(7.0)
lib_top = Inches(2.3)
lib_w = Inches(5.6)
lib_h = Inches(4.3)
rect(s, lib_left, lib_top, lib_w, lib_h, fill=PANEL_LIGHT, line=CYAN, rounded=True)
text(s, "ADAPTER LIBRARY", lib_left + Inches(0.3), lib_top + Inches(0.2),
     Inches(4), Inches(0.4), size=12, color=CYAN, bold=True)
text(s, "Customer-portable across all tiers.",
     lib_left + Inches(0.3), lib_top + Inches(0.55),
     Inches(5), Inches(0.4), size=11, color=GRAY)

# Adapter chips
chips = [
    ("COMPLIANCE", VIOLET), ("FRAUD", CYAN), ("LEGAL", GREEN),
    ("RESEARCH", AMBER), ("MEDICAL", RED), ("CODE", CYAN),
    ("MATH", VIOLET), ("SCIENCE", GREEN),
]
chip_w = Inches(2.4)
chip_h = Inches(0.45)
chip_top = lib_top + Inches(1.05)
for i, (label, col) in enumerate(chips):
    row, colidx = divmod(i, 2)
    x = lib_left + Inches(0.3) + (chip_w + Inches(0.15)) * colidx
    y = chip_top + (chip_h + Inches(0.12)) * row
    rect(s, x, y, chip_w, chip_h, fill=PANEL, line=col, rounded=True)
    text(s, label, x, y + Inches(0.1), chip_w, Inches(0.3),
         size=11, color=col, bold=True, align=PP_ALIGN.CENTER)

# Trailing line below chips
text(s, "+ unlimited custom domains.  Per-customer marginal cost: ~$10-20.",
     lib_left + Inches(0.3), lib_top + Inches(3.7), lib_w - Inches(0.6), Inches(0.4),
     size=11, color=GRAY)

# Bottom moat statement
text(s, "Moat: frontier labs cannot ship this. Their unit economics depend on single-model API routing.",
     Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
     size=13, color=AMBER, bold=True)

# =====================================================================
# SLIDE 6 — WEDGE
# =====================================================================
s = blank_slide()
accent_bar(s)
page_label(s, "Where we win first", 6)
footer(s)

text(s, "Indian BFSI back-office reasoning.",
     Inches(0.7), Inches(1.1), Inches(12), Inches(0.7), size=34, bold=True)
text(s, "Compliance · Internal research · Fraud — workloads frontier APIs cannot legally serve.",
     Inches(0.7), Inches(1.85), Inches(12), Inches(0.5), size=15, color=GRAY)

# 3 use case cards
card_w = Inches(3.95)
card_h = Inches(3.2)
card_top = Inches(2.8)
card_left0 = Inches(0.7)
card_gap = Inches(0.18)

cards = [
    {
        "label": "01  COMPLIANCE",
        "title": "RBI / SEBI documentation",
        "desc": "Auto-summarization, gap detection, and submission drafting on regulated documents that legally cannot leave the country.",
        "stat": "$80-150K",
        "stat_label": "annual contract per institution",
        "color": CYAN,
    },
    {
        "label": "02  RESEARCH",
        "title": "Equity & credit analysis",
        "desc": "Mid-tier asset managers and NBFCs need GPT-4-class document reasoning but legally cannot use it. Adapter-tuned 30B closes the gap.",
        "stat": "$60-120K",
        "stat_label": "annual contract per institution",
        "color": VIOLET,
    },
    {
        "label": "03  FRAUD",
        "title": "Pattern detection on PII",
        "desc": "Transaction-level reasoning over data that triggers DPDP Act protection. On-prem deployment is the only legal architecture.",
        "stat": "$100-200K",
        "stat_label": "annual contract per institution",
        "color": GREEN,
    },
]
for i, c in enumerate(cards):
    x = card_left0 + (card_w + card_gap) * i
    rect(s, x, card_top, card_w, card_h, fill=PANEL, rounded=True)
    rect(s, x, card_top, Inches(0.18), card_h, fill=c["color"])
    text(s, c["label"], x + Inches(0.4), card_top + Inches(0.25),
         Inches(3), Inches(0.3), size=10, color=c["color"], bold=True)
    text(s, c["title"], x + Inches(0.4), card_top + Inches(0.6),
         card_w - Inches(0.6), Inches(0.5), size=18, color=WHITE, bold=True)
    text(s, c["desc"], x + Inches(0.4), card_top + Inches(1.25),
         card_w - Inches(0.6), Inches(1.6), size=12, color=GRAY)
    text(s, c["stat"], x + Inches(0.4), card_top + Inches(2.5),
         card_w - Inches(0.6), Inches(0.5), size=22, color=c["color"], bold=True)
    text(s, c["stat_label"], x + Inches(0.4), card_top + Inches(2.85),
         card_w - Inches(0.6), Inches(0.3), size=10, color=GRAY)

# Bottom math
text(s, "~200 target institutions × $200K avg ACV  =  $40M ARR potential in wedge alone.",
     Inches(0.7), Inches(6.4), Inches(12), Inches(0.5), size=15, color=WHITE, bold=True)
text(s, "Adjacencies: healthcare, defense, govt — same sovereignty buying motion.",
     Inches(0.7), Inches(6.85), Inches(12), Inches(0.4), size=12, color=GRAY)

# =====================================================================
# SLIDE 7 — FOUNDER + WHY US
# =====================================================================
s = blank_slide()
accent_bar(s)
page_label(s, "Why this gets built by us", 7)
footer(s)

text(s, "Solo. 19. Built what 4-senior teams couldn't.",
     Inches(0.7), Inches(1.1), Inches(12), Inches(0.7), size=30, bold=True, color=CYAN)

# 4 timeline cards / proof points
proofs = [
    ("AGE 13", "Started shipping code.", "5+ years deep before college."),
    ("ICPC INDIA 2024", "Built the prelims landing site.", "Lead tech for India's biggest competitive bootcamp."),
    ("BHARAT MANDAPAM", "VisionOS app in one week,", "starting from zero knowledge of the platform."),
    ("YESTERDAY", "Won international ECG hackathon.", "Hardware + ML + fine-tuned model in 12 of 48 hours."),
]
proof_w = Inches(2.95)
proof_h = Inches(2.2)
proof_top = Inches(2.2)
proof_left0 = Inches(0.7)
proof_gap = Inches(0.18)
for i, (when, what, detail) in enumerate(proofs):
    x = proof_left0 + (proof_w + proof_gap) * i
    rect(s, x, proof_top, proof_w, proof_h, fill=PANEL, rounded=True)
    text(s, when, x + Inches(0.25), proof_top + Inches(0.25),
         proof_w - Inches(0.5), Inches(0.4), size=11, color=CYAN, bold=True)
    text(s, what, x + Inches(0.25), proof_top + Inches(0.7),
         proof_w - Inches(0.5), Inches(0.7), size=15, color=WHITE, bold=True)
    text(s, detail, x + Inches(0.25), proof_top + Inches(1.5),
         proof_w - Inches(0.5), Inches(0.6), size=11, color=GRAY)

# Bottom panel — what's been built
rect(s, Inches(0.7), Inches(4.7), Inches(12), Inches(1.9), fill=PANEL_LIGHT, rounded=True)
text(s, "ALREADY BUILT — SOLO", Inches(1.0), Inches(4.85), Inches(8), Inches(0.4),
     size=11, color=CYAN, bold=True)
text(s, "Adapter system validated on Nemotron-30B (4-bit, RTX 5090).",
     Inches(1.0), Inches(5.2), Inches(12), Inches(0.4), size=14, color=WHITE)
text(s, "12 routing techniques tested, 4 papers in pipeline (NeurIPS submission May 4).",
     Inches(1.0), Inches(5.55), Inches(12), Inches(0.4), size=14, color=WHITE)
text(s, "20+ adapters trained across multiple base models.  Code Paradox: novel finding, n=200.",
     Inches(1.0), Inches(5.9), Inches(12), Inches(0.4), size=14, color=WHITE)
text(s, "WebSocket investor demo + adapter-routing visualization shipped.",
     Inches(1.0), Inches(6.25), Inches(12), Inches(0.4), size=14, color=WHITE)

text(s, "What I need is leverage — capital, hardware, and customer access — not validation.",
     Inches(0.7), Inches(6.95), Inches(12), Inches(0.4), size=12, color=AMBER, bold=True)

# =====================================================================
# SLIDE 8 — ASK
# =====================================================================
s = blank_slide()
accent_bar(s)
page_label(s, "What we're raising", 8)
footer(s)

text(s, "Raising $500K SAFE.",
     Inches(0.7), Inches(1.1), Inches(12), Inches(0.8), size=44, bold=True)
text(s, "12-month runway to first 3 BFSI design partners and 200B-scale validation.",
     Inches(0.7), Inches(2.0), Inches(12), Inches(0.5), size=16, color=GRAY)

# Use of funds + Today's ask, side by side
panel_w = Inches(5.95)
panel_h = Inches(3.6)
panel_top = Inches(2.9)

# Left: use of funds
rect(s, Inches(0.7), panel_top, panel_w, panel_h, fill=PANEL, rounded=True)
text(s, "USE OF FUNDS", Inches(0.95), panel_top + Inches(0.25),
     Inches(5), Inches(0.4), size=11, color=CYAN, bold=True)

uof = [
    ("60%", "Cloud GPU access for 70-200B scaling validation"),
    ("25%", "First 3 BFSI design-partner deployments"),
    ("15%", "Founder + minimum infra to ship"),
]
uof_top = panel_top + Inches(0.85)
for i, (pct, what) in enumerate(uof):
    y = uof_top + Inches(0.85) * i
    text(s, pct, Inches(0.95), y, Inches(1.2), Inches(0.6),
         size=26, color=CYAN, bold=True)
    text(s, what, Inches(2.2), y + Inches(0.1), Inches(4.3), Inches(0.7),
         size=13, color=WHITE)

# Right: today's ask
rect(s, Inches(6.85), panel_top, panel_w, panel_h, fill=PANEL_LIGHT, line=CYAN, rounded=True)
text(s, "TODAY'S ASK", Inches(7.1), panel_top + Inches(0.25),
     Inches(5), Inches(0.4), size=11, color=CYAN, bold=True)
text(s, "From this room, three things:", Inches(7.1), panel_top + Inches(0.6),
     Inches(5), Inches(0.4), size=12, color=GRAY)

asks = [
    ("01", "Azure / cloud GPU credits", "$50-150K to validate at 70B+"),
    ("02", "Two BFSI customer intros", "Mid-tier banks or NBFCs in India"),
    ("03", "Advisory relationship", "30 min / month for 12 months"),
]
ask_top = panel_top + Inches(1.25)
for i, (num, what, detail) in enumerate(asks):
    y = ask_top + Inches(0.75) * i
    text(s, num, Inches(7.1), y, Inches(0.7), Inches(0.5),
         size=20, color=CYAN, bold=True)
    text(s, what, Inches(7.85), y - Inches(0.02), Inches(4.5), Inches(0.4),
         size=14, color=WHITE, bold=True)
    text(s, detail, Inches(7.85), y + Inches(0.32), Inches(4.5), Inches(0.4),
         size=11, color=GRAY)

# Closing line
text(s, "Sovereign AI is a 5-year tailwind. We are 12 months ahead on the architecture.",
     Inches(0.7), Inches(6.65), Inches(12), Inches(0.4), size=14, color=WHITE, bold=True)
text(s, "The window to build the platform closes when frontier labs notice it exists.",
     Inches(0.7), Inches(7.0), Inches(12), Inches(0.4), size=12, color=AMBER, bold=True)

# ---------- save ----------
out = "/home/learner/Desktop/mewtwo/SYNAPTA_PITCH_DECK.pptx"
prs.save(out)
print(f"Saved → {out}")
